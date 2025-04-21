from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import re
import traceback
import docx.oxml.shared
from docx.oxml.ns import qn
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn
import lxml.etree as ET


class DocumentGenerator:
    @staticmethod
    def _format_slide_content(slide, section_title, content, main_color, accent_color, text_color):
        """优化的内容排版方法，使PPT布局更合理美观"""
        try:
            # 先设置标题
            slide_title = slide.shapes.title
            if slide_title:
                slide_title.text = section_title
                for paragraph in slide_title.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(36)
                        run.font.bold = True
                        run.font.color.rgb = main_color
            
            # 智能分析内容，提取关键点
            content_lines = content.split('\n')
            content_lines = [line.strip() for line in content_lines if line.strip()]
            
            # 如果内容少于3行，增加字体大小并居中显示
            if len(content_lines) <= 3:
                left = Inches(1.5)
                top = Inches(2.0)
                width = Inches(10)
                height = Inches(4)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                for i, line in enumerate(content_lines):
                    if i > 0:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = line
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.size = Pt(28)
                        run.font.bold = True
                        run.font.color.rgb = accent_color
                return
            
            # 如果有4-10行内容，使用单栏布局，但增加突出显示
            if len(content_lines) <= 10:
                left = Inches(1.0)
                top = Inches(1.8)
                width = Inches(11)
                height = Inches(5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                for i, line in enumerate(content_lines):
                    if i > 0:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    
                    # 如果行以"-"或"•"开头，保留原样，否则添加"•"
                    if not line.startswith(('-', '•', '·', '*')):
                        p.text = f"• {line}"
                    else:
                        p.text = line
                    
                    # 对重要内容进行突出显示
                    is_important = any(keyword in line.lower() for keyword in ['重要', '关键', '核心', '优势', '主要', '特点'])
                    
                    for run in p.runs:
                        if is_important:
                            run.font.size = Pt(24)
                            run.font.bold = True
                            run.font.color.rgb = accent_color
                        else:
                            run.font.size = Pt(22)
                            run.font.color.rgb = text_color
                return
            
            # 对于内容较多的情况，使用双栏布局
            # 计算每栏的行数
            half_point = len(content_lines) // 2
            left_content = content_lines[:half_point]
            right_content = content_lines[half_point:]
            
            # 左栏
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.5), Inches(5))
            left_tf = left_box.text_frame
            left_tf.word_wrap = True
            
            for i, line in enumerate(left_content):
                if i > 0:
                    p = left_tf.add_paragraph()
                else:
                    p = left_tf.paragraphs[0]
                
                # 如果行以"-"或"•"开头，保留原样，否则添加"•"
                if not line.startswith(('-', '•', '·', '*')):
                    p.text = f"• {line}"
                else:
                    p.text = line
                
                # 对重要内容进行突出显示
                is_important = any(keyword in line.lower() for keyword in ['重要', '关键', '核心', '优势', '主要', '特点'])
                
                for run in p.runs:
                    if is_important:
                        run.font.size = Pt(22)
                        run.font.bold = True
                        run.font.color.rgb = accent_color
                    else:
                        run.font.size = Pt(20)
                        run.font.color.rgb = text_color
            
            # 右栏
            right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.7), Inches(5.5), Inches(5))
            right_tf = right_box.text_frame
            right_tf.word_wrap = True
            
            for i, line in enumerate(right_content):
                if i > 0:
                    p = right_tf.add_paragraph()
                else:
                    p = right_tf.paragraphs[0]
                
                # 如果行以"-"或"•"开头，保留原样，否则添加"•"
                if not line.startswith(('-', '•', '·', '*')):
                    p.text = f"• {line}"
                else:
                    p.text = line
                
                # 对重要内容进行突出显示
                is_important = any(keyword in line.lower() for keyword in ['重要', '关键', '核心', '优势', '主要', '特点'])
                
                for run in p.runs:
                    if is_important:
                        run.font.size = Pt(22)
                        run.font.bold = True
                        run.font.color.rgb = accent_color
                    else:
                        run.font.size = Pt(20)
                        run.font.color.rgb = text_color
                        
        except Exception as e:
            print(f"优化幻灯片内容排版失败 - {e}")
            # 简单备用方案
            try:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(5))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = content[:1000] # 截取前1000个字符，避免内容过多
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = text_color
            except:
                pass

    @staticmethod
    def _create_detail_slide(ppt, section_title, content, page_index, total_pages, current_slide_count, main_color, accent_color, text_color):
        """创建内容详细幻灯片，使用优化布局"""
        # 创建详细内容页
        detail_layout = ppt.slide_layouts[5] if len(ppt.slide_layouts) > 5 else ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(detail_layout)
        
        # 设置标题
        try:
            slide_title = slide.shapes.title
            if slide_title:
                if total_pages > 1:
                    slide_title.text = f"{section_title} ({page_index+1}/{total_pages})"
                else:
                    slide_title.text = section_title
                
                # 设置标题格式
                for paragraph in slide_title.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(36)
                        run.font.bold = True
                        run.font.color.rgb = main_color
        except Exception as e:
            print(f"设置标题失败 - {e}")
        
        # 使用优化的内容排版
        DocumentGenerator._format_slide_content(slide, section_title, content, main_color, accent_color, text_color)
        
        return slide

    @staticmethod
    def generate_ppt(title, outline_data, content_data=None, page_limit=None):
        """生成PPT演示文稿"""
        # 创建PPT演示文稿
        ppt = Presentation()
        
        # 设置幻灯片大小为16:9（更现代的比例）
        ppt.slide_width = Inches(13.33)
        ppt.slide_height = Inches(7.5)
        
        # 定义颜色主题
        main_color = RGBColor(39, 71, 125)  # 深蓝色
        accent_color = RGBColor(79, 129, 189)  # 亮蓝色
        title_color = RGBColor(0, 0, 0)  # 黑色
        subtitle_color = RGBColor(64, 64, 64)  # 深灰色
        text_color = RGBColor(89, 89, 89)  # 中灰色
        
        # 1. 创建标题页
        title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title_slide.shapes.title.text = title
        
        # 设置标题页字体更大，更突出
        for paragraph in title_slide.shapes.title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(54)  # 增大字体
                run.font.bold = True
                run.font.color.rgb = title_color
        
        # 如果有副标题占位符，设置一个简单的副标题
        if hasattr(title_slide, "placeholders") and len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = "专业报告"
            for paragraph in subtitle.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(32)  # 增大副标题字体
                    run.font.color.rgb = subtitle_color
        
        # 2. 创建目录页
        toc_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        toc_slide.shapes.title.text = "目录"
        
        # 设置目录标题格式
        for paragraph in toc_slide.shapes.title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.color.rgb = main_color
                run.font.bold = True
        
        # 添加目录内容
        content_shape = None
        for shape in toc_slide.shapes:
            if shape != toc_slide.shapes.title and shape.has_text_frame:
                content_shape = shape
                break
        
        if not content_shape:
            left = Inches(2)
            top = Inches(2)
            width = Inches(9)
            height = Inches(4)
            content_shape = toc_slide.shapes.add_textbox(left, top, width, height)
            
        # 填充目录内容
        tf = content_shape.text_frame
        tf.clear()  # 清除所有现有文本
        
        for i, section in enumerate(outline_data):
            p = tf.add_paragraph()
            p.text = f"{i+1}. {section['title']}"
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.color.rgb = accent_color
                run.font.bold = True
        
        # 跟踪当前幻灯片数量
        current_slide_count = 2  # 已经有标题页和目录页
        
        # 计算分配每个章节的幻灯片数量（确保每个章节都有内容）
        if page_limit is not None:
            # 可用幻灯片数量（减去标题、目录和结束页）
            available_slides = max(len(outline_data), page_limit - 3)
            
            # 确保每个章节至少有一页用于概述
            min_slides_per_section = 1
            
            # 计算每个章节可以分配的额外幻灯片数
            if len(outline_data) * min_slides_per_section < available_slides:
                extra_slides = available_slides - (len(outline_data) * min_slides_per_section)
                slides_per_section = {section["title"]: min_slides_per_section for section in outline_data}
                
                # 分配额外幻灯片
                for i in range(extra_slides):
                    # 循环分配给各章节
                    section_title = outline_data[i % len(outline_data)]["title"]
                    slides_per_section[section_title] += 1
            else:
                # 如果幻灯片不够，每章节只分配一页
                slides_per_section = {section["title"]: min_slides_per_section for section in outline_data}
            
            print(f"幻灯片分配计划: {slides_per_section}")
        else:
            slides_per_section = None  # 不限制
        
        # 3. 为每个章节创建内容幻灯片
        for section_index, section in enumerate(outline_data):
            section_title = section["title"]
            section_slide_count = 0  # 跟踪每个章节的幻灯片数
            
            # 检查是否已达到总页数限制
            if page_limit is not None and current_slide_count >= page_limit:
                print(f"警告：达到总页数限制({page_limit})，后续章节将使用简化模式")
                
                # 如果已经达到页数限制，但还有未处理的章节，采用简化模式
                # 每个剩余章节只生成一个概述页
                for remaining_section in outline_data[section_index:]:
                    if current_slide_count >= page_limit:
                        print(f"达到严格页数限制，无法为{remaining_section['title']}创建页面")
                        break
                        
                    # 创建简化的章节概述页
                    overview_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                    current_slide_count += 1
                    
                    # 设置标题
                    overview_slide.shapes.title.text = remaining_section["title"]
                    for paragraph in overview_slide.shapes.title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = main_color
                            run.font.bold = True
                    
                    # 获取内容形状
                    content_shape = None
                    for shape in overview_slide.shapes:
                        if shape != overview_slide.shapes.title and shape.has_text_frame:
                            content_shape = shape
                            break
                    
                    if not content_shape:
                        try:
                            left = Inches(0.5)
                            top = Inches(1.5)
                            width = Inches(12)
                            height = Inches(5)
                            content_shape = overview_slide.shapes.add_textbox(left, top, width, height)
                        except Exception as e:
                            print(f"创建内容文本框失败 - {e}")
                            continue
                    
                    # 显示大纲点
                    if content_shape:
                        try:
                            tf = content_shape.text_frame
                            tf.clear()
                            
                            # 设置章节概述文本
                            p = tf.add_paragraph()
                            p.text = "章节要点"
                            p.alignment = PP_ALIGN.LEFT
                            for run in p.runs:
                                run.font.size = Pt(28)
                                run.font.color.rgb = accent_color
                                run.font.bold = True
                            
                            # 添加空行
                            tf.add_paragraph()
                            
                            # 添加大纲点
                            for point in remaining_section['content']:
                                p = tf.add_paragraph()
                                p.text = f"• {point}"
                                p.level = 0
                                p.alignment = PP_ALIGN.LEFT
                                for run in p.runs:
                                    run.font.size = Pt(24)
                                    run.font.color.rgb = text_color
                        except Exception as e:
                            print(f"设置概要内容失败 - {e}")
                
                # 添加结束页
                if current_slide_count < page_limit:
                    DocumentGenerator._add_ending_slide(ppt, main_color, accent_color)
                
                return ppt
            
            # 本章节最大可用幻灯片数
            max_section_slides = None
            if slides_per_section:
                max_section_slides = slides_per_section.get(section_title, 1)
                print(f"章节 '{section_title}' 分配的幻灯片数: {max_section_slides}")
            
            # 3.1 创建章节概述页
            overview_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            current_slide_count += 1
            section_slide_count += 1
            
            # 设置标题
            overview_slide.shapes.title.text = section_title
            for paragraph in overview_slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(36)
                    run.font.color.rgb = main_color
                    run.font.bold = True
            
            # 准备概述页内容
            overview_content = f"章节概述:\n\n"
            for point in section['content']:
                overview_content += f"• {point}\n"
                
            # 使用优化的内容排版方法
            DocumentGenerator._format_slide_content(overview_slide, section_title, overview_content, main_color, accent_color, text_color)
            
            # 获取章节的详细内容（如果有）
            if content_data and section_title in content_data and content_data[section_title].strip():
                # 使用AI生成的内容
                detailed_content = content_data[section_title].strip()
                
                # 记录内容获取成功
                print(f"成功获取章节 '{section_title}' 的详细内容 ({len(detailed_content)} 字符)")
                
                # 检查是否达到本章节幻灯片限制
                if max_section_slides is not None and section_slide_count >= max_section_slides:
                    print(f"已达到章节'{section_title}'的幻灯片限制({max_section_slides})，跳过详细内容")
                    continue
                
                # 检查内容是否非常少（少于100字符），如果是则不使用分页而是放在一页
                if len(detailed_content) < 100:
                    # 内容很少，不需要创建额外的页面
                    print(f"章节'{section_title}'内容很少(小于100字符)，不创建额外页面")
                    continue
                
                # 将内容分段显示
                if '\n\n' in detailed_content:
                    paragraphs = detailed_content.split('\n\n')
                else:
                    # 单换行的情况，简单按行分割
                    paragraphs = detailed_content.split('\n')
                
                # 过滤空段落并确保至少有一个段落
                paragraphs = [p for p in paragraphs if p.strip()]
                if not paragraphs:
                    paragraphs = [detailed_content]
                
                # 根据页面限制分配内容
                if max_section_slides is not None:
                    # 计算还可以创建的详细内容页数
                    remaining_slides = max_section_slides - section_slide_count
                    
                    if remaining_slides <= 0:
                        print(f"章节'{section_title}'没有剩余幻灯片额度，跳过详细内容")
                        continue
                        
                    # 如果段落太多，需要合并
                    if len(paragraphs) > remaining_slides:
                        paragraphs_per_slide = max(1, len(paragraphs) // remaining_slides + 1)
                        merged_paragraphs = []
                        
                        for i in range(0, len(paragraphs), paragraphs_per_slide):
                            merged = '\n\n'.join(paragraphs[i:i+paragraphs_per_slide])
                            merged_paragraphs.append(merged)
                        
                        paragraphs = merged_paragraphs[:remaining_slides]
                        print(f"合并内容：将{len(paragraphs)}个段落合并到{len(merged_paragraphs)}页")
                
                # 为每个段落创建一页（如果还有可用页数）
                for i, paragraph_content in enumerate(paragraphs):
                    # 检查是否达到章节页数限制
                    if max_section_slides is not None and section_slide_count >= max_section_slides:
                        print(f"达到章节'{section_title}'的幻灯片限制({max_section_slides})，停止添加详细内容")
                        break
                    
                    # 检查是否达到总页数限制
                    if page_limit is not None and current_slide_count >= page_limit:
                        print(f"达到总页数限制({page_limit})，停止添加详细内容")
                        break
                    
                    # 创建详细内容页并增加计数
                    slide = DocumentGenerator._create_detail_slide(
                        ppt, section_title, paragraph_content, i, len(paragraphs),
                        current_slide_count, main_color, accent_color, text_color
                    )
                    current_slide_count += 1
                    section_slide_count += 1
                    
            else:
                # 内容数据不存在
                print(f"章节 '{section_title}' 没有详细内容数据，仅显示大纲要点")
        
        # 4. 如果页数未达到限制，创建结束页
        if page_limit is None or current_slide_count < page_limit:
            DocumentGenerator._add_ending_slide(ppt, main_color, accent_color)
        
        return ppt
    
    @staticmethod
    def _add_ending_slide(ppt, main_color, accent_color):
        """添加结束页"""
        end_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        
        # 设置结束页标题
        title_shape = end_slide.shapes.title
        title_shape.text = "谢谢观看"
        
        # 设置结束页副标题
        if hasattr(end_slide, "placeholders") and len(end_slide.placeholders) > 1:
            subtitle = end_slide.placeholders[1]
            subtitle.text = "欢迎提问与讨论"
        
        # 设置结束页的格式
        for shape in end_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        if shape == title_shape:
                            run.font.size = Pt(44)
                            run.font.color.rgb = main_color
                            run.font.bold = True
                        else:
                            run.font.size = Pt(28)
                            run.font.color.rgb = accent_color
    
    @staticmethod
    def generate_word(title, outline_data, content_data=None):
        """生成Word文档
        
        Args:
            title (str): 文档标题
            outline_data (list): 大纲数据，格式为[{'title': 章节标题, 'content': [内容列表]}]
            content_data (dict, optional): 详细内容数据，格式为{'章节标题': '详细内容文本'}
            
        Returns:
            Document: 生成的Word文档对象
        """
        # 导入必要的模块
        import docx.oxml.shared
        from docx.oxml.ns import qn
        
        doc = Document()
        
        # 定义中文字体，确保兼容性
        chinese_font = '微软雅黑'
        fallback_fonts = ['宋体', 'SimSun', 'Microsoft YaHei', 'SimHei', 'Arial Unicode MS']
        
        # 全局应用字体的健壮版本
        def apply_chinese_font(run):
            # 设置直接的字体属性
            run.font.name = chinese_font
            
            try:
                # 获取或创建rPr元素
                rPr = run._element.get_or_add_rPr()
                
                # 设置东亚语言
                try:
                    lang = rPr.xpath('./w:lang')
                    if not lang:
                        lang = docx.oxml.shared.OxmlElement('w:lang')
                        lang.set(qn('w:eastAsia'), 'zh-CN')
                        rPr.append(lang)
                    else:
                        lang[0].set(qn('w:eastAsia'), 'zh-CN')
                except Exception as e:
                    print(f"设置语言时出错: {e}")
                
                # 设置rFonts (字体)
                try:
                    rFonts = rPr.xpath('./w:rFonts')
                    if not rFonts:
                        rFonts = docx.oxml.shared.OxmlElement('w:rFonts')
                        rFonts.set(qn('w:eastAsia'), chinese_font)
                        rPr.append(rFonts)
                    else:
                        rFonts[0].set(qn('w:eastAsia'), chinese_font)
                except Exception as e:
                    print(f"设置字体时出错: {e}")
                    
                # 确保此文本使用东亚字体
                try:
                    # 确保此文本使用东亚字体
                    hint = docx.oxml.shared.OxmlElement('w:eastAsianLayout')
                    hint.set(qn('w:id'), '1')
                    rPr.append(hint)
                except Exception:
                    pass
                    
            except Exception as e:
                print(f"应用中文字体失败: {e}")
        
        # 设置一个helper函数来处理段落左右对齐
        def set_paragraph_alignment(paragraph, alignment=WD_ALIGN_PARAGRAPH.JUSTIFY):
            paragraph.alignment = alignment
            # 直接通过XML设置来强制对齐方式
            p = paragraph._p
            pPr = p.get_or_add_pPr()
            jc = docx.oxml.shared.OxmlElement('w:jc')
            jc.set(qn('w:val'), 'both' if alignment == WD_ALIGN_PARAGRAPH.JUSTIFY else 'left')
            for old_jc in pPr.findall(qn('w:jc')):
                pPr.remove(old_jc)
            pPr.append(jc)
        
        # 设置文档默认样式
        style_normal = doc.styles['Normal']
        style_normal.font.name = chinese_font
        style_normal.font.size = DocxPt(10.5)  # 设置正文默认字号为五号(10.5磅)
        
        # 设置标题样式
        for i in range(1, 10):
            style_name = f'Heading {i}'
            if style_name in doc.styles:
                style = doc.styles[style_name]
                style.font.name = chinese_font
                
                # 使用XML设置东亚字体
                rPr = style._element.get_or_add_rPr()
                if rPr is not None:
                    rFonts = rPr.get_or_add_rFonts()
                    rFonts.set(qn('w:eastAsia'), chinese_font)
                
                # 设置标题字号
                heading_sizes = {
                    1: 22,    # 一级标题: 22磅
                    2: 16,    # 二级标题: 16磅
                    3: 15,    # 三级标题: 15磅
                    4: 14,    # 四级标题: 14磅
                    5: 13,    # 五级标题: 13磅
                    6: 12,    # 六级标题: 12磅
                    7: 11,    # 七级标题: 11磅
                    8: 10.5,  # 八级标题: 10.5磅
                    9: 10.5,  # 九级标题: 10.5磅
                }
                if i in heading_sizes:
                    style.font.size = DocxPt(heading_sizes[i])
        
        # 添加标题
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            apply_chinese_font(run)
            run.font.size = DocxPt(22)  # 设置文档标题字号
            run.font.bold = True
        
        # 添加目录标题
        toc_heading = doc.add_heading("目录", 1)
        set_paragraph_alignment(toc_heading, WD_ALIGN_PARAGRAPH.LEFT)
        for run in toc_heading.runs:
            apply_chinese_font(run)
        
        # 添加目录内容
        for i, section in enumerate(outline_data, 1):
            para = doc.add_paragraph()
            set_paragraph_alignment(para, WD_ALIGN_PARAGRAPH.LEFT)
            run = para.add_run(f"{i}. {section['title']}")
            apply_chinese_font(run)
            run.bold = True
        
        # 添加分页符
        doc.add_page_break()
        
        # 添加内容
        for i, section in enumerate(outline_data, 1):
            section_title = section['title']
            
            # 添加章节标题
            heading = doc.add_heading(f"{i}. {section_title}", 1)
            set_paragraph_alignment(heading, WD_ALIGN_PARAGRAPH.LEFT)
            for run in heading.runs:
                apply_chinese_font(run)
            
            # 添加章节概要（大纲点）
            summary_para = doc.add_paragraph()
            set_paragraph_alignment(summary_para, WD_ALIGN_PARAGRAPH.LEFT)
            summary_run = summary_para.add_run("章节概要:")
            apply_chinese_font(summary_run)
            summary_run.bold = True
            
            for point in section['content']:
                para = doc.add_paragraph(style='List Bullet')
                set_paragraph_alignment(para, WD_ALIGN_PARAGRAPH.JUSTIFY)
                point_run = para.add_run(point)
                apply_chinese_font(point_run)
            
            # 添加详细内容（如果有）
            if content_data and section_title in content_data:
                # 添加分隔
                doc.add_paragraph()
                
                # 添加详细内容标题
                detail_para = doc.add_paragraph()
                set_paragraph_alignment(detail_para, WD_ALIGN_PARAGRAPH.LEFT)
                detail_run = detail_para.add_run("详细内容:")
                apply_chinese_font(detail_run)
                detail_run.bold = True
                
                # 添加AI生成的内容，处理可能的Markdown格式
                detailed_content = content_data[section_title]
                
                # 处理内容，移除可能的Markdown标记
                clean_content = detailed_content
                # 替换常见Markdown标记为纯文本
                clean_content = re.sub(r'#{1,6}\s+', '', clean_content)  # 移除标题标记
                clean_content = re.sub(r'\*\*(.*?)\*\*', r'\1', clean_content)  # 移除粗体标记
                clean_content = re.sub(r'\*(.*?)\*', r'\1', clean_content)  # 移除斜体标记
                clean_content = re.sub(r'`(.*?)`', r'\1', clean_content)  # 移除代码标记
                
                # 分段处理内容
                paragraphs = clean_content.split('\n\n')
                if len(paragraphs) <= 1:  # 如果没有足够的段落分隔
                    paragraphs = clean_content.split('\n')  # 尝试按单个换行符分割
                
                for para_text in paragraphs:
                    if para_text.strip():
                        # 检查是否为列表项
                        if para_text.strip().startswith(('-', '*', '•')):
                            # 添加为项目符号
                            bullet_para = doc.add_paragraph(style='List Bullet')
                            set_paragraph_alignment(bullet_para, WD_ALIGN_PARAGRAPH.JUSTIFY)
                            # 移除前导符号
                            clean_text = re.sub(r'^[-*•]\s+', '', para_text.strip())
                            bullet_run = bullet_para.add_run(clean_text)
                            apply_chinese_font(bullet_run)
                        else:
                            # 普通段落
                            normal_para = doc.add_paragraph()
                            set_paragraph_alignment(normal_para, WD_ALIGN_PARAGRAPH.JUSTIFY)
                            normal_run = normal_para.add_run(para_text.strip())
                            apply_chinese_font(normal_run)
        
            # 每个章节后添加分页符，除非是最后一个章节
            if i < len(outline_data):
                doc.add_page_break()
        
        # 设置文档页面格式为A4
        section = doc.sections[0]
        section.page_height = DocxPt(841.9)  # A4高度
        section.page_width = DocxPt(595.3)   # A4宽度
        
        # 设置页边距(单位是磅)
        section.top_margin = DocxPt(72)     # 上边距1英寸
        section.bottom_margin = DocxPt(72)  # 下边距1英寸
        section.left_margin = DocxPt(72)    # 左边距1英寸
        section.right_margin = DocxPt(72)   # 右边距1英寸
        
        # 设置RTL为禁用
        try:
            sectPr = section._sectPr
            if sectPr is not None:
                bidi = docx.oxml.shared.OxmlElement('w:bidi')
                bidi.set(qn('w:val'), '0')
                sectPr.append(bidi)
        except Exception as e:
            print(f"设置RTL失败: {e}")
        
        return doc 